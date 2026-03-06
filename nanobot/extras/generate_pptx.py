#!/usr/bin/env python3
"""
Gerador de apresentações PowerPoint (.pptx)
Uso: python generate_pptx.py [arquivo_saida.pptx]

Dependências: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import sys
import os

# === CORES PADRÃO ===
AZUL_ESCURO = RGBColor(31, 78, 121)
AZUL_CLARO = RGBColor(68, 114, 196)
BRANCO = RGBColor(255, 255, 255)
CINZA = RGBColor(89, 89, 89)

def criar_apresentacao(caminho_saida="out/apresentacao.pptx"):
    """Cria uma apresentação PowerPoint de exemplo"""

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # === SLIDE 1: TÍTULO ===
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide1 = prs.slides.add_slide(slide_layout)

    # Fundo
    background = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = AZUL_ESCURO
    background.line.fill.background()

    # Título
    titulo = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = titulo.text_frame
    p = tf.paragraphs[0]
    p.text = "Título da Apresentação"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
    tf2 = subtitulo.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Subtítulo ou nome do apresentador"
    p2.font.size = Pt(28)
    p2.font.color.rgb = BRANCO
    p2.alignment = PP_ALIGN.CENTER

    # === SLIDE 2: CONTEÚDO COM BULLETS ===
    slide_layout2 = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(slide_layout2)

    # Título do slide
    titulo2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
    tf = titulo2.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Linha decorativa
    linha = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    linha.fill.solid()
    linha.fill.fore_color.rgb = AZUL_CLARO
    linha.line.fill.background()

    # Bullets
    bullets = slide2.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11), Inches(5))
    tf = bullets.text_frame
    tf.word_wrap = True

    itens = [
        "Introdução e contexto",
        "Análise dos dados",
        "Resultados obtidos",
        "Próximos passos",
        "Conclusão"
    ]

    for i, item in enumerate(itens):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = CINZA
        p.space_after = Pt(18)

    # === SLIDE 3: DOIS COLUNAS ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    titulo3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
    tf = titulo3.text_frame
    p = tf.paragraphs[0]
    p.text = "Comparativo"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Coluna esquerda
    col1 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tf = col1.text_frame
    p = tf.paragraphs[0]
    p.text = "Antes"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_CLARO
    for texto in ["Processo manual", "Alto custo", "Erros frequentes"]:
        p = tf.add_paragraph()
        p.text = f"• {texto}"
        p.font.size = Pt(20)
        p.font.color.rgb = CINZA

    # Coluna direita
    col2 = slide3.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
    tf = col2.text_frame
    p = tf.paragraphs[0]
    p.text = "Depois"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_CLARO
    for texto in ["Processo automatizado", "Redução de custos", "Alta precisão"]:
        p = tf.add_paragraph()
        p.text = f"• {texto}"
        p.font.size = Pt(20)
        p.font.color.rgb = CINZA

    # === SLIDE 4: GRÁFICO ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    titulo4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
    tf = titulo4.text_frame
    p = tf.paragraphs[0]
    p.text = "Resultados"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Gráfico de barras
    chart_data = CategoryChartData()
    chart_data.categories = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai', 'Jun']
    chart_data.add_series('2023', (150, 180, 220, 195, 250, 280))
    chart_data.add_series('2024', (180, 210, 260, 240, 300, 350))

    x, y, cx, cy = Inches(1), Inches(1.5), Inches(11), Inches(5.5)
    chart = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True

    # === SLIDE 5: CONCLUSÃO ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    bg = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_ESCURO
    bg.line.fill.background()

    # Texto central
    conclusao = slide5.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(2))
    tf = conclusao.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Dúvidas?"
    p2.font.size = Pt(32)
    p2.font.color.rgb = BRANCO
    p2.alignment = PP_ALIGN.CENTER

    # Criar diretório
    os.makedirs(os.path.dirname(caminho_saida) if os.path.dirname(caminho_saida) else '.', exist_ok=True)

    prs.save(caminho_saida)
    print(f"Apresentação criada: {caminho_saida}")
    return caminho_saida


# === FUNÇÕES AUXILIARES PARA USO COM LLMs ===

def nova_apresentacao(widescreen=True):
    """Cria nova apresentação"""
    prs = Presentation()
    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs

def adicionar_slide_titulo(prs, titulo, subtitulo="", cor_fundo=None):
    """Adiciona slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if cor_fundo:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = cor_fundo
        bg.line.fill.background()

    # Título
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRANCO if cor_fundo else AZUL_ESCURO
    p.alignment = PP_ALIGN.CENTER

    if subtitulo:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitulo
        p2.font.size = Pt(28)
        p2.font.color.rgb = BRANCO if cor_fundo else CINZA
        p2.alignment = PP_ALIGN.CENTER

    return slide

def adicionar_slide_conteudo(prs, titulo, bullets):
    """Adiciona slide com título e lista de bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Título
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Bullets
    bx = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11), Inches(5.5))
    tf = bx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = CINZA
        p.space_after = Pt(14)

    return slide

def adicionar_slide_grafico(prs, titulo, categorias, series_dados):
    """
    Adiciona slide com gráfico de barras

    Args:
        series_dados: dict {'Nome Serie': [valores]}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    chart_data = CategoryChartData()
    chart_data.categories = categorias
    for nome, valores in series_dados.items():
        chart_data.add_series(nome, valores)

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(11), Inches(5.5),
        chart_data
    )

    return slide

def salvar_apresentacao(prs, caminho):
    """Salva a apresentação"""
    os.makedirs(os.path.dirname(caminho) if os.path.dirname(caminho) else '.', exist_ok=True)
    prs.save(caminho)
    print(f"Apresentação salva: {caminho}")


if __name__ == "__main__":
    saida = sys.argv[1] if len(sys.argv) > 1 else "out/apresentacao_exemplo.pptx"
    criar_apresentacao(saida)
