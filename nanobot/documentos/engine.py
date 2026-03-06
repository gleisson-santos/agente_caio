"""
Document Generation Engine — creates professional documents with real content.

Supports: DOCX, PDF, PPTX, XLSX
Each function accepts a title + content dict and produces a file.
"""

from __future__ import annotations

import os
import re
from datetime import datetime
from pathlib import Path
from loguru import logger

OUT_DIR = Path(os.getcwd()) / "out"


def ensure_out() -> Path:
    OUT_DIR.mkdir(exist_ok=True)
    return OUT_DIR


def _ts() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ── DOCX ─────────────────────────────────────────────────────────────

def create_docx(title: str, sections: list[dict], filename: str | None = None) -> str:
    """
    Create a professional Word document.
    sections: [{"heading": "...", "content": "..."}, ...]
    """
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    t = doc.add_heading(title, level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadata
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(f"Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')} — Caio Corp")
    run.font.size = Pt(9)
    run.font.italic = True

    doc.add_paragraph("")  # spacer

    for sec in sections:
        heading = sec.get("heading", "")
        content = sec.get("content", "")

        if heading:
            doc.add_heading(heading, level=1)

        for para in content.split("\n"):
            para = para.strip()
            if not para:
                continue
            if para.startswith("- ") or para.startswith("• "):
                doc.add_paragraph(para[2:], style="List Bullet")
            elif para.startswith("| "):
                # Simple table row — skip for now, add as paragraph
                doc.add_paragraph(para)
            else:
                doc.add_paragraph(para)

    out = ensure_out()
    fname = filename or f"DOC_{_ts()}.docx"
    path = out / fname
    doc.save(str(path))
    logger.info("Engine: DOCX created → {}", fname)
    return fname


# ── PDF ──────────────────────────────────────────────────────────────

def create_pdf(title: str, sections: list[dict], filename: str | None = None) -> str:
    """Create a professional PDF report."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors

    out = ensure_out()
    fname = filename or f"PDF_{_ts()}.pdf"
    path = out / fname

    doc = SimpleDocTemplate(str(path), pagesize=A4,
                            topMargin=2 * cm, bottomMargin=2 * cm,
                            leftMargin=2.5 * cm, rightMargin=2.5 * cm)

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="DocTitle", fontSize=20, leading=26,
                               alignment=1, spaceAfter=12,
                               textColor=HexColor("#1e3a5f")))
    styles.add(ParagraphStyle(name="SectionHead", fontSize=14, leading=18,
                               spaceAfter=8, spaceBefore=16,
                               textColor=HexColor("#2563eb")))
    styles.add(ParagraphStyle(name="Body", fontSize=11, leading=15,
                               spaceAfter=6))
    styles.add(ParagraphStyle(name="Meta", fontSize=9, leading=12,
                               alignment=1, spaceAfter=20,
                               textColor=HexColor("#6b7280")))

    story = []
    story.append(Paragraph(title, styles["DocTitle"]))
    story.append(Paragraph(
        f"Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')} — Caio Corp",
        styles["Meta"]
    ))
    story.append(Spacer(1, 12))

    for sec in sections:
        heading = sec.get("heading", "")
        content = sec.get("content", "")
        if heading:
            story.append(Paragraph(heading, styles["SectionHead"]))
            
        lines = content.split("\n")
        table_data = []

        for para in lines:
            para = para.strip()
            
            # Detect Markdown tables
            if "|" in para:
                # Ignore markdown separator lines like |---|---|
                if re.match(r'^[\s\|:-]+$', para):
                    continue
                
                # Split by | and clean up cells
                row = [cell.strip() for cell in para.split("|")[1:-1] if cell.strip() or para.startswith("|")]
                if not row: # Fallback if split fails or no boundary pipes
                    row = [cell.strip() for cell in para.split("|") if cell.strip()]
                
                if row:
                    # Apply paragraph styling to cells to allow wrapping
                    formatted_row = []
                    for cell in row:
                        cell_clean = re.sub(r'[*_]{1,2}', '', cell) # remove bold/italic markup
                        cell_clean = cell_clean.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        formatted_row.append(Paragraph(cell_clean, styles["Body"]))
                    table_data.append(formatted_row)
                continue
            else:
                # If we were building a table and now hit normal text, render the table
                if table_data:
                    t = Table(table_data, colWidths=[(doc.width / len(table_data[0]))] * len(table_data[0]))
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), HexColor("#2563eb")),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), HexColor("#f8fafc")),
                        ('GRID', (0, 0), (-1, -1), 0.5, HexColor("#cbd5e1"))
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 10))
                    table_data = []

            # Normal Text Processing
            # Clean markup for PDF
            para = re.sub(r'[*_]{1,2}', '', para)
            if para:
                # Escape XML chars
                para = para.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(para, styles["Body"]))
                
        # Flush any remaining table at the end of section
        if table_data:
            t = Table(table_data, colWidths=[(doc.width / len(table_data[0]))] * len(table_data[0]))
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), HexColor("#2563eb")),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), HexColor("#f8fafc")),
                ('GRID', (0, 0), (-1, -1), 0.5, HexColor("#cbd5e1"))
            ]))
            story.append(t)
            story.append(Spacer(1, 10))

        story.append(Spacer(1, 6))

    doc.build(story)
    logger.info("Engine: PDF created → {}", fname)
    return fname


# ── PPTX ─────────────────────────────────────────────────────────────

def create_pptx(title: str, slides: list[dict], filename: str | None = None) -> str:
    """
    Create a professional PowerPoint presentation.
    slides: [{"title": "...", "content": "..."}, ...]
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"Caio Corp — {datetime.now().strftime('%d/%m/%Y')}"

    # Content slides
    for s in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = s.get("title", "")
        
        # Parse content for tables
        lines = s.get("content", "").split("\n")
        table_rows = []
        normal_lines = []
        
        for line in lines:
            line = line.strip()
            if "|" in line:
                if re.match(r'^[\s\|:-]+$', line):
                    continue
                row = [cell.strip() for cell in line.split("|")[1:-1] if cell.strip() or line.startswith("|")]
                if not row:
                    row = [cell.strip() for cell in line.split("|") if cell.strip()]
                if row:
                    table_rows.append(row)
            else:
                if line:
                    normal_lines.append(line)
        
        if table_rows:
            # Render Table
            # Remove the default text placeholder to make room for the table
            sp = slide.shapes[1].element
            sp.getparent().remove(sp)
            
            rows = len(table_rows)
            cols = max(len(row) for row in table_rows)
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(11.3)
            height = Inches(0.8 * rows)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            for i, row in enumerate(table_rows):
                for j, cell in enumerate(row):
                    if j < cols:
                        table.cell(i, j).text = re.sub(r'[*_]{1,2}', '', cell) # clear markdown
                        for paragraph in table.cell(i, j).text_frame.paragraphs:
                            paragraph.font.size = Pt(14)
            
            # If there's extra text, add it above or below (simplified: just log it or add textbox)
            if normal_lines:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11.3), Inches(0.7))
                tf = txBox.text_frame
                tf.text = " ".join([re.sub(r'[*_]{1,2}', '', l) for l in normal_lines])
                
        else:
            # Render Normal Text
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for line in normal_lines:
                # Clean markup for PPTX
                line = re.sub(r'[*_]{1,2}', '', line)
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)

    out = ensure_out()
    fname = filename or f"PPTX_{_ts()}.pptx"
    path = out / fname
    prs.save(str(path))
    logger.info("Engine: PPTX created → {}", fname)
    return fname


# ── XLSX ─────────────────────────────────────────────────────────────

def create_xlsx(title: str, sheets: list[dict], filename: str | None = None) -> str:
    """
    Create a professional Excel workbook.
    sheets: [{"name": "...", "headers": [...], "rows": [[...], ...]}, ...]
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()

    for idx, sheet_data in enumerate(sheets):
        if idx == 0:
            ws = wb.active
            ws.title = sheet_data.get("name", "Dados")
        else:
            ws = wb.create_sheet(title=sheet_data.get("name", f"Sheet{idx + 1}"))

        headers = sheet_data.get("headers", [])
        rows = sheet_data.get("rows", [])

        # Title row
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 1))
        cell = ws.cell(row=1, column=1, value=title)
        cell.font = Font(bold=True, size=14, color="1e3a5f")
        cell.alignment = Alignment(horizontal="center")

        # Header row
        header_fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin")
        )

        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = thin_border
            ws.column_dimensions[chr(64 + col_idx)].width = max(len(str(header)) + 4, 15)

        # Data rows
        for row_idx, row in enumerate(rows, 4):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center" if isinstance(value, (int, float)) else "left")

    out = ensure_out()
    fname = filename or f"XLSX_{_ts()}.xlsx"
    path = out / fname
    wb.save(str(path))
    logger.info("Engine: XLSX created → {}", fname)
    return fname


# ── SAMPLE CONTENT (for Quick Actions) ───────────────────────────────

SAMPLE_RELATORIO = {
    "title": "Relatório Gerencial — Caio Corp",
    "sections": [
        {"heading": "1. Sumário Executivo", "content": "Este relatório apresenta os indicadores operacionais e financeiros do período.\nOs resultados demonstram crescimento consistente em todas as áreas monitoradas.\nDestaque: aumento de 23% na eficiência operacional comparado ao período anterior."},
        {"heading": "2. Indicadores de Desempenho (KPIs)", "content": "- Taxa de resolução de chamados: 94.2%\n- Tempo médio de resposta: 2h 15min\n- Satisfação do cliente (NPS): 78 pontos\n- Uptime dos sistemas: 99.7%\n- Projetos entregues no prazo: 87%"},
        {"heading": "3. Análise Financeira", "content": "Receita acumulada: R$ 245.800,00\nCustos operacionais: R$ 128.400,00\nMargem operacional: 47.7%\nInvestimentos em infraestrutura: R$ 32.000,00\nROI estimado: 18.5%"},
        {"heading": "4. Ações e Próximos Passos", "content": "- Implementar automação de processos repetitivos (Q2)\n- Expandir capacidade de atendimento em 30%\n- Migrar infraestrutura para cloud híbrida\n- Treinar equipe em novas ferramentas de IA"},
        {"heading": "5. Conclusão", "content": "Os resultados indicam uma trajetória positiva. A equipe está alinhada com as metas estabelecidas e os investimentos em tecnologia estão gerando retorno mensurável. Recomenda-se manter o foco em automação e eficiência operacional para o próximo trimestre."},
    ]
}

SAMPLE_APRESENTACAO = {
    "title": "Apresentação Executiva — Caio Corp",
    "slides": [
        {"title": "Visão Geral do Período", "content": "Crescimento de 23% na eficiência operacional\nReceita: R$ 245.800,00 (+15% vs período anterior)\n6 projetos entregues com sucesso\n94.2% taxa de resolução de chamados"},
        {"title": "Indicadores Financeiros", "content": "Receita acumulada: R$ 245.800\nCustos operacionais: R$ 128.400\nMargem operacional: 47.7%\nROI: 18.5%"},
        {"title": "Conquistas do Período", "content": "Implementação do sistema de IA para atendimento\nRedução de 40% no tempo de resposta\nNPS de 78 pontos (meta: 75)\nUptime de 99.7% nos sistemas críticos"},
        {"title": "Desafios Identificados", "content": "Necessidade de escalar a infraestrutura\nIntegração com novos fornecedores\nTreinamento da equipe em ferramentas avançadas\nOtimização de custos de cloud"},
        {"title": "Plano de Ação Q2", "content": "Automação de processos repetitivos\nExpansão da capacidade em 30%\nMigração para cloud híbrida\nPrograma de capacitação em IA"},
    ]
}

SAMPLE_PLANILHA = {
    "title": "Controle Operacional — Caio Corp",
    "sheets": [
        {
            "name": "Indicadores",
            "headers": ["Indicador", "Meta", "Realizado", "Status", "%"],
            "rows": [
                ["Taxa de Resolução", "90%", "94.2%", "✅ Atingido", 94.2],
                ["Tempo de Resposta", "3h", "2h 15min", "✅ Atingido", 100],
                ["NPS", "75", "78", "✅ Atingido", 104],
                ["Uptime", "99.5%", "99.7%", "✅ Atingido", 100.2],
                ["Projetos no Prazo", "85%", "87%", "✅ Atingido", 102.4],
                ["Custo por Chamado", "R$ 45", "R$ 38", "✅ Atingido", 115.8],
            ]
        },
        {
            "name": "Financeiro",
            "headers": ["Item", "Jan", "Fev", "Mar", "Total"],
            "rows": [
                ["Receita", 78500, 82300, 85000, 245800],
                ["Custos", 41200, 43100, 44100, 128400],
                ["Margem", 37300, 39200, 40900, 117400],
                ["Investimentos", 10000, 12000, 10000, 32000],
            ]
        }
    ]
}

SAMPLE_CONTRATO = {
    "title": "Contrato de Prestação de Serviços",
    "sections": [
        {"heading": "CONTRATO DE PRESTAÇÃO DE SERVIÇOS DE CONSULTORIA", "content": "Pelo presente instrumento particular de contrato de prestação de serviços, as partes abaixo qualificadas têm entre si, justo e acertado, o presente contrato."},
        {"heading": "CLÁUSULA PRIMEIRA — DAS PARTES", "content": "CONTRATANTE: [Razão Social da Empresa Contratante], inscrita no CNPJ sob o nº [XX.XXX.XXX/0001-XX], com sede na [endereço completo], representada por [nome do representante].\n\nCONTRATADA: Caio Corp Tecnologia Ltda, especialista em soluções de inteligência artificial e automação."},
        {"heading": "CLÁUSULA SEGUNDA — DO OBJETO", "content": "O presente contrato tem por objeto a prestação de serviços de consultoria em tecnologia da informação, incluindo:\n- Análise e otimização de processos\n- Implementação de soluções de IA\n- Suporte técnico especializado\n- Treinamento da equipe"},
        {"heading": "CLÁUSULA TERCEIRA — DO PRAZO", "content": "O presente contrato terá vigência de 12 (doze) meses, contados a partir da data de sua assinatura, podendo ser renovado por iguais períodos mediante acordo entre as partes."},
        {"heading": "CLÁUSULA QUARTA — DO VALOR", "content": "Pelo serviço descrito, a CONTRATANTE pagará à CONTRATADA o valor mensal de R$ [valor], totalizando R$ [valor total] para o período contratual. O pagamento será efetuado até o dia 10 de cada mês."},
        {"heading": "CLÁUSULA QUINTA — DAS OBRIGAÇÕES", "content": "Ambas as partes se comprometem a cumprir fielmente as obrigações aqui estabelecidas, sob pena de rescisão contratual e aplicação das sanções cabíveis."},
        {"heading": "ASSINATURAS", "content": f"Local e Data: ________________, {datetime.now().strftime('%d de %B de %Y')}\n\n\n___________________________\nCONTRATANTE\n\n\n___________________________\nCONTRATADA"},
    ]
}


# ── Quick Generate (with defaults) ───────────────────────────────────

def quick_generate(doc_type: str) -> str:
    """Generate a document with professional sample content."""
    if doc_type == "pdf":
        return create_pdf(SAMPLE_RELATORIO["title"], SAMPLE_RELATORIO["sections"])
    elif doc_type == "pptx":
        return create_pptx(SAMPLE_APRESENTACAO["title"], SAMPLE_APRESENTACAO["slides"])
    elif doc_type == "xlsx":
        return create_xlsx(SAMPLE_PLANILHA["title"], SAMPLE_PLANILHA["sheets"])
    elif doc_type == "docx":
        return create_docx(SAMPLE_CONTRATO["title"], SAMPLE_CONTRATO["sections"])
    else:
        raise ValueError(f"Unsupported type: {doc_type}")


def generate_from_ai_content(content: str, doc_format: str, title: str = "Documento IA") -> str:
    """Parse AI-generated text into structured sections and create a file."""
    sections = []
    current_heading = ""
    current_content = []

    import re
    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("# ") or stripped.startswith("## "):
            if current_heading or current_content:
                sections.append({"heading": current_heading, "content": "\n".join(current_content)})
                current_content = []
            current_heading = stripped.lstrip("#").strip()
        elif stripped.startswith("### "):
            if current_heading or current_content:
                sections.append({"heading": current_heading, "content": "\n".join(current_content)})
                current_content = []
            current_heading = stripped.lstrip("#").strip()
        else:
            current_content.append(stripped)

    if current_heading or current_content:
        sections.append({"heading": current_heading, "content": "\n".join(current_content)})

    # If no sections parsed, treat entire content as one section
    if not sections:
        sections = [{"heading": "", "content": content}]

    if doc_format == "pdf":
        return create_pdf(title, sections)
    elif doc_format == "pptx":
        # Convert sections to slides
        slides = [{"title": s["heading"] or f"Slide {i+1}", "content": s["content"]} for i, s in enumerate(sections)]
        return create_pptx(title, slides)
    elif doc_format == "xlsx":
        # Convert to simple data sheet
        rows = []
        for s in sections:
            if s["heading"]:
                rows.append([s["heading"], ""])
            for line in s["content"].split("\n"):
                line = line.strip()
                if line:
                    rows.append(["", line])
        sheets = [{"name": "Conteúdo", "headers": ["Seção", "Detalhe"], "rows": rows}]
        return create_xlsx(title, sheets)
    else:  # docx
        return create_docx(title, sections)

TEMPLATES_DIR = Path(__file__).parent / "templates"

def list_generated_documents(out_path: str | None = None) -> list[dict[str, Any]]:
    """List all files in the output directory with metadata."""
    out = Path(out_path) if out_path else OUT_DIR
    if not out.exists():
        return []
    
    docs = []
    for f in out.iterdir():
        if f.is_file() and not f.name.startswith("."):
            stats = f.stat()
            docs.append({
                "name": f.name,
                "type": f.suffix.lstrip(".").lower(),
                "size": stats.st_size,
                "sizeFormatted": f"{stats.st_size / 1024:.1f} KB",
                "createdAt": stats.st_mtime,
                "icon": "📄" # default
            })
    # Sort by creation time desc
    docs.sort(key=lambda x: x["createdAt"], reverse=True)
    return docs


def get_templates() -> list[dict[str, Any]]:
    """List available templates with metadata."""
    templates = []
    template_meta = {
        "contrato_prestacao.md": {
            "title": "Contrato de Prestação",
            "description": "Contrato padrão para serviços de consultoria e TI.",
            "icon": "⚖️", "category": "legal",
        },
        "relatorio_gerencial.md": {
            "title": "Relatório Gerencial",
            "description": "Resumo executivo com KPIs e análise de resultados.",
            "icon": "📈", "category": "empresarial",
        },
        "ata_reuniao.md": {
            "title": "Ata de Reunião",
            "description": "Registro formal de decisões e ações acordadas.",
            "icon": "📝", "category": "empresarial",
        },
        "procuracao.md": {
            "title": "Procuração Particular",
            "description": "Modelo para representação legal simples.",
            "icon": "📜", "category": "legal",
        },
        "proposta_comercial.md": {
            "title": "Proposta Comercial",
            "description": "Estrutura profissional para vendas e novos negócios.",
            "icon": "📋", "category": "empresarial",
        },
        "plano_projeto.md": {
            "title": "Plano de Projeto",
            "description": "Cronograma, metas e riscos de execução.",
            "icon": "🗓️", "category": "empresarial",
        },
    }
    
    if TEMPLATES_DIR.exists():
        for f in TEMPLATES_DIR.iterdir():
            if f.is_file() and f.suffix == ".md":
                meta = template_meta.get(f.name, {
                    "title": f.stem.replace("_", " ").title(),
                    "description": "", "icon": "📝", "category": "outro",
                })
                templates.append({
                    "id": f.stem,
                    "filename": f.name,
                    **meta,
                })
    return templates


def get_template_content(template_id: str) -> dict[str, Any]:
    """Get the raw content of a specific template."""
    path = TEMPLATES_DIR / f"{template_id}.md"
    if not path.exists():
        raise FileNotFoundError(f"Template {template_id} not found")
    
    content = path.read_text(encoding="utf-8", errors="replace")
    return {"id": template_id, "content": content}
